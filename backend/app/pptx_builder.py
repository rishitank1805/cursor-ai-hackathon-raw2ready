"""Build PowerPoint from Manus content with varied layouts. Images from Picsum only (no Unsplash)."""

import io
from typing import Any

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Theme colors
COLOR_DARK_PANEL = RGBColor(0x0F, 0x1E, 0x3C)
COLOR_ACCENT = RGBColor(0xE8, 0x6C, 0x25)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT = RGBColor(0xE8, 0xE8, 0xE8)
COLOR_BODY = RGBColor(0xCC, 0xCC, 0xCC)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF8, 0xFC)
COLOR_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)


def _fetch_url(url: str, timeout: float = 10.0) -> bytes | None:
    try:
        with httpx.Client(timeout=timeout, follow_redirects=True) as client:
            r = client.get(url)
            r.raise_for_status()
            return r.content
    except Exception:
        return None


def _resolve_slide_image(slide_data: dict, idx: int, presentation_title: str) -> bytes | None:
    """Image for slide from Picsum only (no Unsplash). Seed from content so each slide/topic differs."""
    query = (slide_data.get("image_search_query") or "").strip()
    if not query:
        title = (slide_data.get("title") or "").strip()
        query = f"{presentation_title} {title}".strip() or "business"
    seed = abs(hash(query + str(idx))) % (10**8)
    url = f"https://picsum.photos/seed/{seed}/800/450"
    return _fetch_url(url)


def _add_textbox(slide, left, top, width, height, text, font_size, bold, color):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = "Calibri"
    p.font.color.rgb = color
    return box


def _add_bullets(slide, left, top, width, height, content_list, color, font_size=16):
    if not content_list:
        return
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(content_list):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.name = "Calibri"
        p.font.color.rgb = color
        p.space_before = Pt(6)
        p.level = 0


# ---- Layout 0: Title slide – business name + full-width image top, subtitle below ----
def _layout_title(slide, slide_data, img_bytes, presentation_title, generated_tagline, business_name: str | None = None):
    title_text = (business_name or slide_data.get("title") or presentation_title).strip()
    subtitle_text = slide_data.get("subtitle", "") or generated_tagline or ""
    # Image top ~55%
    if img_bytes:
        try:
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Inches(0), Inches(0),
                SLIDE_WIDTH, Inches(4.15),
            )
        except Exception:
            pass
    # Dark strip bottom with title + subtitle
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.15), SLIDE_WIDTH, Inches(3.35))
    strip.fill.solid()
    strip.fill.fore_color.rgb = COLOR_DARK_PANEL
    strip.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.15), SLIDE_WIDTH, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_ACCENT
    accent.line.fill.background()
    _add_textbox(slide, Inches(1), Inches(4.6), SLIDE_WIDTH - Inches(2), Inches(1), title_text, 38, True, COLOR_WHITE)
    if subtitle_text:
        _add_textbox(slide, Inches(1), Inches(5.7), SLIDE_WIDTH - Inches(2), Inches(0.7), subtitle_text, 20, False, COLOR_LIGHT)


# ---- Layout 1: Split – left dark panel, right image ----
def _layout_split_left_text(slide, slide_data, img_bytes, presentation_title, generated_tagline, business_name: str | None = None):
    left_w = Inches(6.7)
    right_start = Inches(6.85)
    title_text = slide_data.get("title", "Untitled")
    subtitle_text = slide_data.get("subtitle", "") or (generated_tagline if slide_data.get("slide_number") == 1 else "")
    content_list = slide_data.get("content", [])
    left_margin = Inches(0.55)
    # Left panel
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), left_w, SLIDE_HEIGHT)
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLOR_DARK_PANEL
    panel.line.fill.background()
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), left_w, Inches(0.08))
    acc.fill.solid()
    acc.fill.fore_color.rgb = COLOR_ACCENT
    acc.line.fill.background()
    _add_textbox(slide, left_margin, Inches(0.65), left_w - Inches(0.6), Inches(1.1), title_text, 34, True, COLOR_WHITE)
    if subtitle_text:
        _add_textbox(slide, left_margin, Inches(1.75), left_w - Inches(0.6), Inches(0.6), subtitle_text, 18, False, COLOR_LIGHT)
    content_top = Inches(2.5) if subtitle_text else Inches(1.85)
    _add_bullets(slide, left_margin, content_top, left_w - Inches(0.6), Inches(4.2), content_list, COLOR_BODY, 16)
    # Right image
    if img_bytes:
        try:
            slide.shapes.add_picture(io.BytesIO(img_bytes), right_start, Inches(0), Inches(6.4), SLIDE_HEIGHT)
        except Exception:
            _placeholder(slide, right_start, Inches(0), Inches(6.4), SLIDE_HEIGHT)
    else:
        _placeholder(slide, right_start, Inches(0), Inches(6.4), SLIDE_HEIGHT)


# ---- Layout 2: Split – right dark panel, left image ----
def _layout_split_right_text(slide, slide_data, img_bytes, presentation_title, generated_tagline, business_name: str | None = None):
    right_w = Inches(6.7)
    left_w_img = Inches(6.6)
    title_text = slide_data.get("title", "Untitled")
    subtitle_text = slide_data.get("subtitle", "")
    content_list = slide_data.get("content", [])
    text_left = Inches(6.65)
    margin = Inches(0.5)
    # Left image
    if img_bytes:
        try:
            slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0), Inches(0), left_w_img, SLIDE_HEIGHT)
        except Exception:
            _placeholder(slide, Inches(0), Inches(0), left_w_img, SLIDE_HEIGHT)
    else:
        _placeholder(slide, Inches(0), Inches(0), left_w_img, SLIDE_HEIGHT)
    # Right panel
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, text_left - Inches(0.05), Inches(0), right_w + Inches(0.1), SLIDE_HEIGHT)
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLOR_DARK_PANEL
    panel.line.fill.background()
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, text_left, Inches(0), right_w, Inches(0.08))
    acc.fill.solid()
    acc.fill.fore_color.rgb = COLOR_ACCENT
    acc.line.fill.background()
    _add_textbox(slide, text_left + margin, Inches(0.65), right_w - Inches(0.6), Inches(1.1), title_text, 34, True, COLOR_WHITE)
    if subtitle_text:
        _add_textbox(slide, text_left + margin, Inches(1.75), right_w - Inches(0.6), Inches(0.6), subtitle_text, 18, False, COLOR_LIGHT)
    content_top = Inches(2.5) if subtitle_text else Inches(1.85)
    _add_bullets(slide, text_left + margin, content_top, right_w - Inches(0.6), Inches(4.2), content_list, COLOR_BODY, 16)


# ---- Layout 3: Image full-width top, text below (light background) ----
def _layout_image_top(slide, slide_data, img_bytes, presentation_title, generated_tagline, business_name: str | None = None):
    title_text = slide_data.get("title", "Untitled")
    subtitle_text = slide_data.get("subtitle", "")
    content_list = slide_data.get("content", [])
    img_h = Inches(3.8)
    # Image
    if img_bytes:
        try:
            slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0), Inches(0), SLIDE_WIDTH, img_h)
        except Exception:
            _placeholder(slide, Inches(0), Inches(0), SLIDE_WIDTH, img_h)
    else:
        _placeholder(slide, Inches(0), Inches(0), SLIDE_WIDTH, img_h)
    # Accent line under image
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), img_h, SLIDE_WIDTH, Inches(0.06))
    acc.fill.solid()
    acc.fill.fore_color.rgb = COLOR_ACCENT
    acc.line.fill.background()
    # Light content area
    content_top = img_h + Inches(0.4)
    _add_textbox(slide, Inches(0.6), content_top, SLIDE_WIDTH - Inches(1.2), Inches(0.9), title_text, 32, True, COLOR_DARK_TEXT)
    if subtitle_text:
        _add_textbox(slide, Inches(0.6), content_top + Inches(0.95), SLIDE_WIDTH - Inches(1.2), Inches(0.5), subtitle_text, 18, False, RGBColor(0x55, 0x55, 0x66))
    bullet_top = content_top + (Inches(1.6) if subtitle_text else Inches(1.0))
    _add_bullets(slide, Inches(0.6), bullet_top, SLIDE_WIDTH - Inches(1.2), Inches(2.2), content_list, COLOR_DARK_TEXT, 18)


# ---- Layout 4: Text above (dark strip), image full-width bottom ----
def _layout_image_bottom(slide, slide_data, img_bytes, presentation_title, generated_tagline, business_name: str | None = None):
    title_text = slide_data.get("title", "Untitled")
    subtitle_text = slide_data.get("subtitle", "")
    content_list = slide_data.get("content", [])
    strip_h = Inches(3.4)
    img_top = strip_h
    # Top dark strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, strip_h)
    strip.fill.solid()
    strip.fill.fore_color.rgb = COLOR_DARK_PANEL
    strip.line.fill.background()
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06))
    acc.fill.solid()
    acc.fill.fore_color.rgb = COLOR_ACCENT
    acc.line.fill.background()
    _add_textbox(slide, Inches(0.6), Inches(0.5), SLIDE_WIDTH - Inches(1.2), Inches(0.9), title_text, 32, True, COLOR_WHITE)
    if subtitle_text:
        _add_textbox(slide, Inches(0.6), Inches(1.45), SLIDE_WIDTH - Inches(1.2), Inches(0.5), subtitle_text, 18, False, COLOR_LIGHT)
    _add_bullets(slide, Inches(0.6), Inches(2.0), SLIDE_WIDTH - Inches(1.2), Inches(1.2), content_list[:3], COLOR_BODY, 14)
    # Bottom image
    if img_bytes:
        try:
            slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0), img_top, SLIDE_WIDTH, SLIDE_HEIGHT - img_top)
        except Exception:
            _placeholder(slide, Inches(0), img_top, SLIDE_WIDTH, SLIDE_HEIGHT - img_top)
    else:
        _placeholder(slide, Inches(0), img_top, SLIDE_WIDTH, SLIDE_HEIGHT - img_top)


# ---- Layout 5: Narrow image strip left, text right (light) ----
def _layout_strip_left(slide, slide_data, img_bytes, presentation_title, generated_tagline, business_name: str | None = None):
    title_text = slide_data.get("title", "Untitled")
    subtitle_text = slide_data.get("subtitle", "")
    content_list = slide_data.get("content", [])
    strip_w = Inches(4.2)
    text_left = strip_w + Inches(0.4)
    text_width = SLIDE_WIDTH - text_left - Inches(0.5)
    # Left image strip
    if img_bytes:
        try:
            slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0), Inches(0), strip_w, SLIDE_HEIGHT)
        except Exception:
            _placeholder(slide, Inches(0), Inches(0), strip_w, SLIDE_HEIGHT)
    else:
        _placeholder(slide, Inches(0), Inches(0), strip_w, SLIDE_HEIGHT)
    # Accent between strip and text
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, strip_w, Inches(0), Inches(0.08), SLIDE_HEIGHT)
    acc.fill.solid()
    acc.fill.fore_color.rgb = COLOR_ACCENT
    acc.line.fill.background()
    # Text
    _add_textbox(slide, text_left, Inches(0.5), text_width, Inches(0.95), title_text, 30, True, COLOR_DARK_TEXT)
    if subtitle_text:
        _add_textbox(slide, text_left, Inches(1.5), text_width, Inches(0.5), subtitle_text, 16, False, RGBColor(0x55, 0x55, 0x66))
    bullet_top = Inches(2.15) if subtitle_text else Inches(1.6)
    _add_bullets(slide, text_left, bullet_top, text_width, Inches(5), content_list, COLOR_DARK_TEXT, 17)


# ---- Layout 6: Centered quote / key message – large title, small image corner ----
def _layout_centered(slide, slide_data, img_bytes, presentation_title, generated_tagline, business_name: str | None = None):
    title_text = slide_data.get("title", "Untitled")
    subtitle_text = slide_data.get("subtitle", "")
    content_list = slide_data.get("content", [])
    # Full dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_DARK_PANEL
    bg.line.fill.background()
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    acc.fill.solid()
    acc.fill.fore_color.rgb = COLOR_ACCENT
    acc.line.fill.background()
    # Centered title (large)
    _add_textbox(slide, Inches(1.5), Inches(2.2), SLIDE_WIDTH - Inches(3), Inches(1.2), title_text, 40, True, COLOR_WHITE)
    if subtitle_text:
        _add_textbox(slide, Inches(1.5), Inches(3.5), SLIDE_WIDTH - Inches(3), Inches(0.6), subtitle_text, 22, False, COLOR_LIGHT)
    if content_list:
        _add_bullets(slide, Inches(2.5), Inches(4.2), SLIDE_WIDTH - Inches(5), Inches(2.2), content_list[:4], COLOR_BODY, 18)
    # Small image bottom-right corner
    if img_bytes:
        try:
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Inches(9.5), Inches(5.2),
                Inches(3.5), Inches(2.0),
            )
        except Exception:
            pass


def _placeholder(slide, left, top, width, height):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x4A)
    s.line.color.rgb = COLOR_ACCENT


# ---- Thank You slide (always last) ----
def _layout_thank_you(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_DARK_PANEL
    bg.line.fill.background()
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    acc.fill.solid()
    acc.fill.fore_color.rgb = COLOR_ACCENT
    acc.line.fill.background()
    _add_textbox(slide, Inches(2), Inches(2.8), SLIDE_WIDTH - Inches(4), Inches(1.2), "Thank You", 52, True, COLOR_WHITE)
    _add_textbox(slide, Inches(2), Inches(4.1), SLIDE_WIDTH - Inches(4), Inches(0.6), "Questions?", 24, False, COLOR_LIGHT)


# Layout roster: 0=title, 1=split L, 2=split R, 3=img top, 4=img bottom, 5=strip L, 6=centered
LAYOUTS = [
    _layout_title,
    _layout_split_left_text,
    _layout_split_right_text,
    _layout_image_top,
    _layout_image_bottom,
    _layout_strip_left,
    _layout_centered,
]


def build_pptx(
    presentation_title: str,
    slides: list[dict[str, Any]],
    generated_tagline: str | None = None,
    include_images: bool = True,
    business_name: str | None = None,
) -> bytes:
    """
    Build PowerPoint from Manus content only. First slide shows business_name; always ends with Thank You.
    Images from Picsum only (no Unsplash).
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]
    n_layouts = len(LAYOUTS)

    for idx, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        img_bytes = None
        if include_images:
            img_bytes = _resolve_slide_image(slide_data, idx, presentation_title)
        layout_idx = 0 if idx == 0 else (1 + (idx - 1) % (n_layouts - 1))
        first_slide_business_name = business_name if idx == 0 else None
        LAYOUTS[layout_idx](slide, slide_data, img_bytes, presentation_title, generated_tagline, first_slide_business_name)
        speaker_notes = slide_data.get("speaker_notes", "")
        if speaker_notes:
            try:
                slide.notes_slide.notes_text_frame.text = speaker_notes
            except Exception:
                pass

    # Always add Thank You slide at the end
    thank_slide = prs.slides.add_slide(blank)
    _layout_thank_you(thank_slide)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def build_pptx_from_response(
    response: dict[str, Any],
    include_images: bool = True,
    business_name: str | None = None,
) -> bytes:
    """Build PPTX from Manus response only. First slide = business name; last slide = Thank You. No Unsplash."""
    slides_data = []
    for s in response.get("slides", []):
        if hasattr(s, "model_dump"):
            slides_data.append(s.model_dump())
        elif isinstance(s, dict):
            slides_data.append(s)
        else:
            slides_data.append({
                "title": getattr(s, "title", "Untitled"),
                "subtitle": getattr(s, "subtitle", None),
                "content": getattr(s, "content", []),
                "speaker_notes": getattr(s, "speaker_notes"),
                "slide_number": getattr(s, "slide_number", len(slides_data) + 1),
                "image_search_query": getattr(s, "image_search_query", None),
            })
    return build_pptx(
        presentation_title=response.get("presentation_title", "Business Pitch"),
        slides=slides_data,
        generated_tagline=response.get("generated_tagline"),
        include_images=include_images,
        business_name=business_name,
    )
